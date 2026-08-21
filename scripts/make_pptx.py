"""Build the Evaluation-1 presentation deck (docs/EVAL1_presentation.pptx).
Run from repo root:  python3 scripts/make_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "docs", "assets")
OUT = os.path.join(ROOT, "docs", "EVAL1_presentation.pptx")

# palette (matches dashboard)
BG = RGBColor(0x0B, 0x0F, 0x1A)
PANEL = RGBColor(0x13, 0x1A, 0x2B)
TEXT = RGBColor(0xE6, 0xEB, 0xF5)
MUTED = RGBColor(0x8B, 0x96, 0xAD)
GREEN = RGBColor(0x3D, 0xDC, 0x84)
BLUE = RGBColor(0x42, 0x85, 0xF4)
AMBER = RGBColor(0xFF, 0xB4, 0x54)
RED = RGBColor(0xFF, 0x5D, 0x5D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def set_para(p, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, space_after=6, bullet=None):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Segoe UI"
    if bullet:
        p.level = 0
        # manual bullet marker
        p.text = (bullet + "  " + text)


def header(slide, kicker, title):
    k = textbox(slide, 0.7, 0.45, 12, 0.4)
    k.text_frame.paragraphs[0].text = kicker.upper()
    for r in k.text_frame.paragraphs[0].runs:
        r.font.size = Pt(13); r.font.color.rgb = GREEN; r.font.bold = True; r.font.name = "Segoe UI"
    t = textbox(slide, 0.7, 0.85, 12, 0.9)
    t.text_frame.paragraphs[0].text = title
    for r in t.text_frame.paragraphs[0].runs:
        r.font.size = Pt(32); r.font.color.rgb = TEXT; r.font.bold = True; r.font.name = "Segoe UI"


def bullets(slide, items, l=0.9, t=2.0, w=11.5, h=5.0):
    tb = textbox(slide, l, t, w, h)
    tf = tb.text_frame
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        marker = it.get("marker", "▸")
        set_para(p, it["text"], size=it.get("size", 20), color=it.get("color", TEXT),
                 bold=it.get("bold", False), space_after=it.get("gap", 12), bullet=marker)
        if it.get("sub"):
            ps = tf.add_paragraph()
            set_para(ps, it["sub"], size=15, color=MUTED, space_after=it.get("gap", 12))
    return tb


# --------------------------------------------------------------------------- #
# Slide 1 — Title
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
bg_path = os.path.join(ASSETS, "title-bg.png")
if os.path.exists(bg_path):
    s.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

tb = textbox(s, 0.9, 2.1, 11.5, 1.4)
p = tb.text_frame.paragraphs[0]
p.text = "AEGIS"
p.alignment = PP_ALIGN.LEFT
for r in p.runs:
    r.font.size = Pt(66); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = "Segoe UI"

tb = textbox(s, 0.9, 3.35, 11.5, 0.9)
p = tb.text_frame.paragraphs[0]
p.text = "Self-Healing Multi-Cloud Platform"
for r in p.runs:
    r.font.size = Pt(28); r.font.color.rgb = GREEN; r.font.bold = True; r.font.name = "Segoe UI"

tb = textbox(s, 0.9, 4.35, 11.5, 1.6)
tf = tb.text_frame
lines = [
    "One application · three cloud free tiers · automatic failover · $0 cost",
    "Final Year Major Project — Evaluation 1",
    "[Your Name]  ·  [College / Department]  ·  August 2026",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.space_after = Pt(6)
    for r in p.runs:
        r.font.size = Pt(16 if i < 2 else 15)
        r.font.color.rgb = MUTED if i < 2 else RGBColor(0xBD, 0xC7, 0xDB)
        r.font.name = "Segoe UI"

# --------------------------------------------------------------------------- #
# Slide 2 — The problem
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Why this project exists", "The Problem")
bullets(s, [
    {"text": "One cloud = one point of failure", "bold": True,
     "sub": "A single provider outage takes your whole application down — you have no control and no recovery path."},
    {"text": "Free tiers are free… but individually unreliable", "bold": True,
     "sub": "Google Cloud Run, Oracle Always Free and Render each give real free hosting — but any one of them can sleep, throttle, or go down."},
    {"text": "Outages cost money and trust", "bold": True,
     "sub": "Downtime means lost users, lost revenue, and damaged reputation — especially for students and early startups with no budget."},
    {"text": "Existing solutions are expensive or locked-in", "bold": True,
     "sub": "Commercial multi-cloud/DR products cost real money and tie you to one vendor's tooling."},
], t=2.0)

# --------------------------------------------------------------------------- #
# Slide 3 — The idea
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Our idea", "Run one app on 3 free tiers — and fail over automatically")
bullets(s, [
    {"text": "Deploy the SAME app to 3 clouds' free tiers", "marker": "1.",
     "sub": "Google Cloud Run · Oracle Always Free · Render — zero cost, no code changes (cloud name via env var)."},
    {"text": "Monitor every cloud every 5 minutes", "marker": "2.",
     "sub": "A health monitor polls /health and stores history in Supabase."},
    {"text": "Route through one smart entry point", "marker": "3.",
     "sub": "A Cloudflare Worker always sends traffic to the healthiest cloud — failover in milliseconds, not minutes."},
    {"text": "Add an intelligence layer on top", "marker": "4.",
     "sub": "Cost savings, carbon-aware routing, policy-as-code, security scanning, and ML-based predictive failover."},
    {"text": "Total running cost: $0", "marker": "5.", "bold": True, "color": GREEN,
     "sub": "Everything runs inside free tiers — the whole point of the project."},
], t=2.0)

# --------------------------------------------------------------------------- #
# Slide 4 — Architecture
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "System design", "Architecture — 5 layers")
arch = os.path.join(ASSETS, "architecture.png")
if os.path.exists(arch):
    s.shapes.add_picture(arch, Inches(3.3), Inches(1.55), height=Inches(5.7))

# --------------------------------------------------------------------------- #
# Slide 5 — What's built (progress for Eval 1)
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Progress — Months 1 to 3", "What's already built")

def status_items():
    return [
        ("✓", GREEN, "Software Requirements Specification (SRS) + architecture diagrams"),
        ("✓", GREEN, "Full repo structure + GitHub Projects board (6-month plan)"),
        ("✓", GREEN, "Core app (Express) with /health contract + Dockerfile"),
        ("✓", GREEN, "2 LIVE cloud deployments — Render + Vercel (both free tiers)"),
        ("✓", GREEN, "https://aegis-platform-pomf.onrender.com/health"),
        ("✓", GREEN, "https://aegis-platform-lyart.vercel.app/health"),
        ("✓", GREEN, "FAILOVER ROUTER LIVE — https://aegis-router.manav32478.workers.dev"),
        ("✓", GREEN, "Health monitor → Supabase (auto-runs every 5 min)"),
        ("✓", GREEN, "Live status page on every cloud (/status)"),
        ("✓", GREEN, "CI/CD pipelines: deploy · monitor · security scan · ML"),
        ("◐", AMBER, "GCP/Oracle optional · ML predictive failover next (Month 4)"),
    ]

tb = textbox(s, 0.9, 1.95, 11.6, 5.2)
tf = tb.text_frame
for i, (mark, color, text) in enumerate(status_items()):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    r1 = p.add_run(); r1.text = f"{mark}  "
    r1.font.size = Pt(20); r1.font.color.rgb = color; r1.font.bold = True; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(20); r2.font.color.rgb = TEXT; r2.font.name = "Segoe UI"

# --------------------------------------------------------------------------- #
# Slide 6 — Live demo
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The demo", "Live: watch failover happen")
bullets(s, [
    {"text": "Open the dashboard", "marker": "▶",
     "sub": "node demo/server.js → localhost:8080 — three clouds, live latency chart."},
    {"text": "Press “Run chaos test”", "marker": "🔥",
     "sub": "A random cloud is forced down for 20 seconds."},
    {"text": "Watch the router fail over", "marker": "▶",
     "sub": "Traffic instantly moves to the next healthiest cloud — the app never goes down."},
    {"text": "Same logic runs in production", "marker": "▶",
     "sub": "The demo router is a direct port of the real Cloudflare Worker (orchestration/router/src/worker.js)."},
], t=2.0)

# --------------------------------------------------------------------------- #
# Slide 7 — Failover proof
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Evidence", "Measured failover from the running system")
fp = os.path.join(ASSETS, "failover-proof.png")
if os.path.exists(fp):
    s.shapes.add_picture(fp, Inches(1.0), Inches(1.7), width=Inches(11.3))

# --------------------------------------------------------------------------- #
# Slide 8 — Roadmap
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Where this goes next", "Roadmap — Months 4 to 6")
bullets(s, [
    {"text": "Month 4 — Intelligence layer", "bold": True, "marker": "▸",
     "sub": "Cost intelligence · carbon-aware routing · policy-as-code (OPA) · security scanning (Trivy/Checkov) · ML predictive failover (Isolation Forest)."},
    {"text": "Month 5 — Product layer", "bold": True, "marker": "▸",
     "sub": "Auth + API keys (Supabase) · REST API with Swagger docs · public status page · load testing with k6."},
    {"text": "Month 6 — Hardening & defense", "bold": True, "marker": "▸",
     "sub": "Unit test suite · final report · demo video · viva preparation."},
    {"text": "Beyond — startup path", "bold": True, "marker": "▸",
     "sub": "The platform + historical data become the first assets of a real product (Section 7 of the blueprint)."},
], t=2.0)

# --------------------------------------------------------------------------- #
# Slide 9 — Why these choices (defense)
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Design decisions", "Why we chose these tools")
rows = [
    ("Decision", "Why"),
    ("Cloudflare Workers (not DNS failover)", "DNS TTL makes failover take minutes; a Worker retries other backends in milliseconds."),
    ("Supabase (Postgres)", "Free hosted database + auth; doubles as ML training data."),
    ("OPA / Rego policies", "Policy decoupled from code — new rules need no app change; auditable."),
    ("Isolation Forest (ML)", "Adapts to each cloud's normal latency instead of one fixed threshold."),
    ("Docker + CI/CD", "Identical image on all 3 clouds; every push is scanned and deployed automatically."),
]
tbl_shape = s.shapes.add_table(len(rows), 2, Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0))
table = tbl_shape.table
table.columns[0].width = Inches(5.4)
table.columns[1].width = Inches(6.1)
for i, (c0, c1) in enumerate(rows):
    for j, txt in enumerate((c0, c1)):
        cell = table.cell(i, j)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if i > 0 else RGBColor(0x1C, 0x2A, 0x4A)
        cell.margin_top = Pt(6); cell.margin_bottom = Pt(6)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        for r in para.runs:
            r.font.size = Pt(15 if i > 0 else 16)
            r.font.color.rgb = TEXT if i > 0 else RGBColor(0x9D, 0xB4, 0xFF)
            r.font.bold = (i == 0) or (j == 0)
            r.font.name = "Segoe UI"

# --------------------------------------------------------------------------- #
# Slide 10 — Thank you
# --------------------------------------------------------------------------- #
s = prs.slides.add_slide(BLANK)
bg(s)
tb = textbox(s, 0.9, 2.6, 11.5, 1.2)
p = tb.text_frame.paragraphs[0]
p.text = "Thank you"
p.alignment = PP_ALIGN.LEFT
for r in p.runs:
    r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = "Segoe UI"
tb = textbox(s, 0.9, 3.9, 11.5, 1.0)
p = tb.text_frame.paragraphs[0]
p.text = "Questions?"
for r in p.runs:
    r.font.size = Pt(26); r.font.color.rgb = GREEN; r.font.name = "Segoe UI"

prs.save(OUT)
print("wrote", OUT)
